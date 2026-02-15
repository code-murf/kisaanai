from pptx import Presentation
import os

pptx_path = "Idea Submission _ AWS AI for Bharat Hackathon.pptx"

if not os.path.exists(pptx_path):
    print(f"Error: file {pptx_path} not found")
    exit(1)

prs = Presentation(pptx_path)

print(f"Presentation has {len(prs.slides)} slides")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text}")
    
    for shape in slide.placeholders:
        print(f"Placeholder index: {shape.placeholder_format.idx}, Name: {shape.name}, Text: {shape.text if hasattr(shape, 'text') else ''}")
    
    for shape in slide.shapes:
        if not shape.is_placeholder and hasattr(shape, "text"):
             print(f"Shape: {shape.name}, Text: {shape.text}")

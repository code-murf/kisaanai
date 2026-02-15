from pptx import Presentation
import sys
import os

def extract_text(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return

    try:
        prs = Presentation(pptx_path)
        print(f"# Content from {os.path.basename(pptx_path)}\n")
        
        for i, slide in enumerate(prs.slides):
            print(f"## Slide {i+1}")
            
            # Extract title
            if slide.shapes.title:
                print(f"### {slide.shapes.title.text}")
            
            # Extract other text
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape == slide.shapes.title:
                    continue
                    
                for paragraph in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in paragraph.runs)
                    if text.strip():
                        print(f"- {text.strip()}")
            print("\n---\n")
            
    except Exception as e:
        print(f"Error processing file: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <path_to_pptx>")
    else:
        extract_text(sys.argv[1])

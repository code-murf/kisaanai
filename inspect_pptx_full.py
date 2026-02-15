from pptx import Presentation
import os

pptx_path = "Idea Submission _ AWS AI for Bharat Hackathon.pptx"
prs = Presentation(pptx_path)

with open("pptx_content.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        
        # Check title
        if slide.shapes.title:
             f.write(f"Title: {slide.shapes.title.text}\n")
        
        # Check all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(f"Shape: {shape.name} | Text: {shape.text}\n")
            if shape.has_text_frame:
                 for paragraph in shape.text_frame.paragraphs:
                     f.write(f"  Paragraph: {paragraph.text}\n")

print("Full content dumped to pptx_content.txt")

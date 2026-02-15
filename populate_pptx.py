from pptx import Presentation
from pptx.util import Inches
import os

pptx_path = "Idea Submission _ AWS AI for Bharat Hackathon.pptx"
output_path = "KisaanAI_Final_Submission.pptx"

if not os.path.exists(pptx_path):
    print(f"Error: file {pptx_path} not found")
    exit(1)

prs = Presentation(pptx_path)

# Content Mapping
content_map = {
    0: { # Slide 1
        "Team Name :": "Team Name : Antigravity AI",
        "Problem Statement :": "Problem Statement : Information Asymmetry & Market Inefficiency in Indian Agriculture",
        "Team Leader Name :": "Team Leader Name : Antigravity Lead"
    },
    1: { # Slide 2
        "Brief about the Idea:": "Brief about the Idea:\n\nKisaanAI is a voice-first, AI-powered platform that empowers farmers with real-time mandi prices, weather forecasts, and crop advisory services in their local language. It bridges the information gap, enabling data-driven decisions for better profitability and reducing distress selling."
    },
    2: { # Slide 3
        "Your solution should be able to explain the following:": "Proposed Solution & USP:\n\n1. Voice-First Interface: Breaks literacy barriers using Bhashini AI, allowing farmers to interact naturally.\n2. Predictive Intelligence: 7-day price forecasts using Temporal Fusion Transformers (TFT) for better market timing.\n3. Smart Routing: Calculates Net Profit (Price - Transport Cost) to recommend the best mandi, not just the nearest one.\n4. Explainable AI: Provides reasoning behind predictions to build trust."
    },
    3: { # Slide 4
        "List of features offered by the solution": "Key Features:\n\n• KisaanCredit (New): AI-based credit scoring using yield forecasts for instant micro-loans.\n• AI Price Forecasting: Accurate predictions for Potato, Onion, Tomato.\n• Voice Assistant: 'Bolo aur Jaano' interface.\n• Smart Routing: Real-time transport cost calculation.\n• WhatsApp Bot: Low-bandwidth access.\n• Crop Doctor: Image-based disease detection."
    },
    4: { # Slide 5
        "Process flow diagram or Use-case diagram": "Process Flow:\n\n1. Farmer asks query via Voice/WhatsApp.\n2. Bhashini translates speech to text.\n3. NLP Engine identifies intent (Price/Weather/Disease).\n4. Backend fetches data from AI Models or Databases.\n5. AI generates personalized response.\n6. Bhashini converts text to speech.\n7. Farmer receives actionable advice."
    },
    5: { # Slide 6
         "Wireframes/Mock diagrams of the proposed solution (optional)": "User Interface Experience:\n\nThe KisaanAI Dashboard features a high-contrast, accessible design:\n- MagicUI Animations: Engaging, smooth interactions.\n- Bento Grid Layout: Clear, modular information display.\n- Responsive Design: Optimized for low-end Android devices and desktops.\n- Visual Cues: Icons and color-coding (Green/Red) for intuitive understanding of trends."
    },
    6: { # Slide 7
        "Architecture diagram of the proposed solution:": "System Architecture:\n\n[Frontend]: Next.js 16 (PWA) + Tailwind CSS + MagicUI\n[Backend]: FastAPI + PostgreSQL (PostGIS) + Redis\n[AI Layer]: \n  - PyTorch (Price Forecasting)\n  - YOLO (Disease Detection)\n  - Bhashini (Voice/Vernacular)\n[Data Pipeline]: Agmarknet (Prices) + Sentinel-2 (Satellite) + IMD (Weather)\n[Infrastructure]: AWS EC2 + S3 + RDS"
    },
    7: { # Slide 8
        "Technologies to be used in the solution:": "Technology Stack:\n\n• Cloud: AWS (EC2, S3, RDS)\n• AI/ML: PyTorch, XGBoost, Temporal Fusion Transformer (TFT)\n• Voice/NLP: Bhashini API, Twilio (WhatsApp)\n• Web: Next.js, Tailwind CSS, Leaflet Maps, MagicUI\n• Backend: FastAPI, PostgreSQL, MinIO"
    },
    8: { # Slide 9
        "Estimated implementation cost (optional):": "Implementation Cost & Feasibility:\n\n• Infrastructure: ~$50/month (Optimized with AWS Spot Instances).\n• Data: Free (Open Government Data Integration).\n• AI Costs: Managed via free tiers for research (Bhashini).\n• Feasibility: MVP ready for deployment; Scalable to 50+ crops."
    },
    9: { # Slide 10
        "Add as per the requirements for the hackathon:": "Business Impact & Future Scope:\n\n• Impact: Projected 15-20% increase in farmer income via smart arbitrage.\n• Scale: Designed for nationwide roll-out with local dialect support.\n• Future: Integration with Farmer Producer Organizations (FPOs) for collective bargaining and direct-to-buyer sales."
    }
}

# Replace Text
for slide_idx, replacements in content_map.items():
    if slide_idx < len(prs.slides):
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for key, value in replacements.items():
                    if key in shape.text:
                        print(f"Replacing text in Slide {slide_idx+1}: {key[:20]}...")
                        shape.text = value

# Insert Dashboard Hero Image into Slide 6 (Index 5 - User Interface)
hero_img = "dashboard_hero.png"
if os.path.exists(hero_img):
    print(f"Inserting image {hero_img} into Slide 6...")
    slide = prs.slides[5] 
    # Adjust placement for Slide 6
    left = Inches(0.5) 
    top = Inches(2.5)
    width = Inches(9) 
    slide.shapes.add_picture(hero_img, left, top, width=width)

# Insert Features Image into Slide 4 (Index 3 - Key Features)
feat_img = "dashboard_features.png"
if os.path.exists(feat_img):
    print(f"Inserting image {feat_img} into Slide 4...")
    slide = prs.slides[3]
    # Adjust placement for Slide 4 (to the right of text potentially, or below)
    left = Inches(5) 
    top = Inches(2.0)
    width = Inches(4.5)
    slide.shapes.add_picture(feat_img, left, top, width=width)

prs.save(output_path)
print(f"Presentation saved to {output_path}")
